"""OfficeAnalyzer — native extraction from DOCX, XLSX, PPTX.

For legacy DOC/XLS/PPT (binary) formats, raises MarkdownConversionError so
the fallback pipeline routes them through MarkItDown or an intermediate
Office→PDF→extract path.
"""

from __future__ import annotations

import logging
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from app.ai_pipeline.utils.md_helpers import rows_to_md
from app.services.markdown_convert_service import MarkdownConversionError

logger = logging.getLogger(__name__)

_MAX_XLSX_ROWS = 200
_MAX_PPTX_SLIDES = 100


def analyze(input_path: str, original_filename: str) -> str:
    """Dispatch to the correct Office analyzer based on extension."""
    ext = Path(original_filename).suffix.lower().lstrip(".")
    logger.debug("OfficeAnalyzer: processing .%s", ext)

    if ext == "docx":
        return _docx(input_path)
    if ext == "xlsx":
        return _xlsx(input_path)
    if ext == "pptx":
        return _pptx(input_path)
    # Legacy binary formats: no native path → let fallback handle them
    raise MarkdownConversionError(
        f"No native analyzer for .{ext}; routing to fallback pipeline."
    )


def _docx(input_path: str) -> str:
    """Extract structured content from DOCX with headings, formatting, and tables."""
    try:
        with zipfile.ZipFile(input_path) as archive:
            xml = archive.read("word/document.xml")
    except (KeyError, zipfile.BadZipFile) as exc:
        raise MarkdownConversionError("DOCX content could not be read.") from exc

    root = ElementTree.fromstring(xml)
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    parts: list[str] = []

    for element in root.findall(".//{%s}body/*" % ns["w"]):
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            md_line = _docx_paragraph(element, ns)
            if md_line:
                parts.append(md_line)
        elif tag == "tbl":
            md_table = _docx_table(element, ns)
            if md_table:
                parts.append(md_table)

    return _require("\n\n".join(parts), "DOCX contains no readable text.")


def _xlsx(input_path: str) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(input_path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[list[str]] = []
        for row in sheet.iter_rows(max_row=_MAX_XLSX_ROWS, values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in values):
                rows.append(values)
        if rows:
            parts.append(f"## {sheet.title}\n\n{rows_to_md(rows)}")
    workbook.close()
    return _require("\n\n".join(parts), "Spreadsheet contains no readable rows.")


def _pptx(input_path: str) -> str:
    from pptx import Presentation

    presentation = Presentation(input_path)
    parts: list[str] = []
    for index, slide in enumerate(
        list(presentation.slides)[:_MAX_PPTX_SLIDES], start=1
    ):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _norm(shape.text)
                if text:
                    texts.append(text)
        if texts:
            parts.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
    return _require("\n\n".join(parts), "Presentation contains no readable text.")


# --- Helpers ---

# Mapping of common Word heading styles to Markdown heading levels
_HEADING_STYLE_MAP = {
    "heading1": 1, "heading 1": 1, "title": 1,
    "heading2": 2, "heading 2": 2, "subtitle": 2,
    "heading3": 3, "heading 3": 3,
    "heading4": 4, "heading 4": 4,
    "heading5": 5, "heading 5": 5,
    "heading6": 6, "heading 6": 6,
}


def _docx_paragraph(para, ns: dict) -> str:
    """Convert a DOCX paragraph to Markdown with heading/list/formatting support."""
    # Detect paragraph style
    style = ""
    ppr = para.find("w:pPr", ns)
    if ppr is not None:
        style_el = ppr.find("w:pStyle", ns)
        if style_el is not None:
            style = style_el.get(f"{{{ns['w']}}}val", "").lower()

        # Check for numbered/bullet list
        num_pr = ppr.find("w:numPr", ns)
        if num_pr is not None:
            text = _get_formatted_runs(para, ns)
            return f"- {text}" if text else ""

    text = _get_formatted_runs(para, ns)
    if not text:
        return ""

    # Map style to heading
    heading_level = _HEADING_STYLE_MAP.get(style, 0)
    if heading_level:
        return f"{'#' * heading_level} {text}"

    return text


def _get_formatted_runs(para, ns: dict) -> str:
    """Extract text from paragraph runs, preserving bold and italic."""
    parts: list[str] = []
    for run in para.findall("w:r", ns):
        text = "".join(t.text or "" for t in run.findall("w:t", ns))
        if not text:
            continue

        # Check formatting properties
        rpr = run.find("w:rPr", ns)
        if rpr is not None:
            is_bold = rpr.find("w:b", ns) is not None
            is_italic = rpr.find("w:i", ns) is not None
            if is_bold and is_italic:
                text = f"***{text}***"
            elif is_bold:
                text = f"**{text}**"
            elif is_italic:
                text = f"*{text}*"

        parts.append(text)

    return "".join(parts).strip()


def _docx_table(tbl, ns: dict) -> str:
    """Extract a DOCX table and convert to Markdown."""
    rows: list[list[str]] = []
    for tr in tbl.findall("w:tr", ns):
        cells: list[str] = []
        for tc in tr.findall("w:tc", ns):
            cell_text = " ".join(
                "".join(t.text or "" for t in p.findall(".//w:t", ns))
                for p in tc.findall("w:p", ns)
            ).strip()
            cells.append(cell_text)
        if cells:
            rows.append(cells)
    if len(rows) >= 2:
        return rows_to_md(rows)
    return ""


def _norm(value: str) -> str:
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def _require(text: str, message: str) -> str:
    if not text.strip():
        raise MarkdownConversionError(message)
    return text

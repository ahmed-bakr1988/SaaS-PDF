"""PDF conversion service — PDF↔PowerPoint, Excel→PDF, PowerPoint→PDF, Sign PDF."""
import os
import io
import logging
import subprocess
import tempfile

logger = logging.getLogger(__name__)


class PDFConvertError(Exception):
    """Custom exception for PDF conversion failures."""
    pass


# ---------------------------------------------------------------------------
# PDF to PowerPoint (PPTX)
# ---------------------------------------------------------------------------
def pdf_to_pptx(input_path: str, output_path: str) -> dict:
    """Convert a PDF to PowerPoint by rendering each page as a slide image.

    Args:
        input_path: Path to the input PDF
        output_path: Path for the output PPTX

    Returns:
        dict with total_slides and output_size

    Raises:
        PDFConvertError: If conversion fails
    """
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches, Emu

        images = convert_from_path(input_path, dpi=200)
        if not images:
            raise PDFConvertError("PDF has no pages or could not be rendered.")

        prs = Presentation()
        # Use widescreen 16:9 layout
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            img_stream = io.BytesIO()
            img.save(img_stream, format="PNG")
            img_stream.seek(0)

            # Scale image to fill slide
            img_w, img_h = img.size
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            ratio = min(slide_w / Emu(int(img_w * 914400 / 200)),
                        slide_h / Emu(int(img_h * 914400 / 200)))
            pic_w = int(img_w * 914400 / 200 * ratio)
            pic_h = int(img_h * 914400 / 200 * ratio)
            left = (slide_w - pic_w) // 2
            top = (slide_h - pic_h) // 2

            slide.shapes.add_picture(img_stream, left, top, pic_w, pic_h)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)

        output_size = os.path.getsize(output_path)
        logger.info(f"PDF→PPTX: {len(images)} slides ({output_size} bytes)")
        return {"total_slides": len(images), "output_size": output_size}

    except PDFConvertError:
        raise
    except Exception as e:
        raise PDFConvertError(f"PDF to PowerPoint conversion failed: {str(e)}")


# ---------------------------------------------------------------------------
# Excel (XLSX) to PDF
# ---------------------------------------------------------------------------
def excel_to_pdf(input_path: str, output_dir: str) -> str:
    """Convert an Excel file to PDF using LibreOffice headless.

    Args:
        input_path: Path to the input XLSX/XLS file
        output_dir: Directory for the output file

    Returns:
        Path to the converted PDF

    Raises:
        PDFConvertError: If conversion fails
    """
    os.makedirs(output_dir, exist_ok=True)
    user_install_dir = tempfile.mkdtemp(prefix="lo_excel2pdf_")

    cmd = [
        "soffice",
        "--headless",
        "--norestore",
        f"-env:UserInstallation=file://{user_install_dir}",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_path,
    ]

    try:
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=120,
            env={**os.environ, "HOME": user_install_dir},
        )

        input_basename = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_dir, f"{input_basename}.pdf")

        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            logger.info(f"Excel→PDF conversion successful: {output_path}")
            return output_path

        if result.returncode != 0:
            stderr = result.stderr or ""
            real_errors = [
                line for line in stderr.strip().splitlines()
                if not line.startswith("Warning: failed to launch javaldx")
            ]
            error_msg = "\n".join(real_errors) if real_errors else stderr
            raise PDFConvertError(f"Conversion failed: {error_msg or 'Unknown error'}")

        raise PDFConvertError("Output file was not created.")

    except subprocess.TimeoutExpired:
        raise PDFConvertError("Conversion timed out. File may be too large.")
    except FileNotFoundError:
        raise PDFConvertError("LibreOffice is not installed on the server.")
    finally:
        import shutil
        shutil.rmtree(user_install_dir, ignore_errors=True)


# ---------------------------------------------------------------------------
# PowerPoint (PPTX) to PDF
# ---------------------------------------------------------------------------
def pptx_to_pdf(input_path: str, output_dir: str) -> str:
    """Convert a PowerPoint file to PDF using LibreOffice headless.

    Args:
        input_path: Path to the input PPTX/PPT file
        output_dir: Directory for the output file

    Returns:
        Path to the converted PDF

    Raises:
        PDFConvertError: If conversion fails
    """
    os.makedirs(output_dir, exist_ok=True)
    user_install_dir = tempfile.mkdtemp(prefix="lo_pptx2pdf_")

    cmd = [
        "soffice",
        "--headless",
        "--norestore",
        f"-env:UserInstallation=file://{user_install_dir}",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_path,
    ]

    try:
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=120,
            env={**os.environ, "HOME": user_install_dir},
        )

        input_basename = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_dir, f"{input_basename}.pdf")

        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            logger.info(f"PPTX→PDF conversion successful: {output_path}")
            return output_path

        if result.returncode != 0:
            stderr = result.stderr or ""
            real_errors = [
                line for line in stderr.strip().splitlines()
                if not line.startswith("Warning: failed to launch javaldx")
            ]
            error_msg = "\n".join(real_errors) if real_errors else stderr
            raise PDFConvertError(f"Conversion failed: {error_msg or 'Unknown error'}")

        raise PDFConvertError("Output file was not created.")

    except subprocess.TimeoutExpired:
        raise PDFConvertError("Conversion timed out. File may be too large.")
    except FileNotFoundError:
        raise PDFConvertError("LibreOffice is not installed on the server.")
    finally:
        import shutil
        shutil.rmtree(user_install_dir, ignore_errors=True)


# ---------------------------------------------------------------------------
# Sign PDF (overlay signature image on a page)
# ---------------------------------------------------------------------------
def sign_pdf(
    input_path: str,
    signature_path: str,
    output_path: str,
    page: int = 0,
    x: float = 100,
    y: float = 100,
    width: float = 200,
    height: float = 80,
) -> dict:
    """Overlay a signature image onto a PDF page.

    Args:
        input_path: Path to the input PDF
        signature_path: Path to the signature image (PNG with transparency)
        output_path: Path for the signed output PDF
        page: 0-based page index to place signature
        x: X coordinate (points from left)
        y: Y coordinate (points from bottom)
        width: Signature width in points
        height: Signature height in points

    Returns:
        dict with total_pages and output_size

    Raises:
        PDFConvertError: If signing fails
    """
    try:
        from PyPDF2 import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.utils import ImageReader

        reader = PdfReader(input_path)
        total_pages = len(reader.pages)
        if total_pages == 0:
            raise PDFConvertError("PDF has no pages.")
        if page < 0 or page >= total_pages:
            raise PDFConvertError(f"Page {page + 1} does not exist (PDF has {total_pages} pages).")

        target_page = reader.pages[page]
        page_box = target_page.mediabox
        page_width = float(page_box.width)
        page_height = float(page_box.height)

        # Create overlay PDF with the signature image
        overlay_stream = io.BytesIO()
        c = rl_canvas.Canvas(overlay_stream, pagesize=(page_width, page_height))
        sig_img = ImageReader(signature_path)
        c.drawImage(sig_img, x, y, width=width, height=height, mask="auto")
        c.save()
        overlay_stream.seek(0)

        overlay_reader = PdfReader(overlay_stream)
        overlay_page = overlay_reader.pages[0]

        writer = PdfWriter()
        for i, pg in enumerate(reader.pages):
            if i == page:
                pg.merge_page(overlay_page)
            writer.add_page(pg)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "wb") as f:
            writer.write(f)

        output_size = os.path.getsize(output_path)
        logger.info(f"Sign PDF: signature on page {page + 1} ({output_size} bytes)")
        return {"total_pages": total_pages, "output_size": output_size, "signed_page": page + 1}

    except PDFConvertError:
        raise
    except Exception as e:
        raise PDFConvertError(f"Failed to sign PDF: {str(e)}")
